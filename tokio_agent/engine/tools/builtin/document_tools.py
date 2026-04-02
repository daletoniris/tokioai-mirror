"""
Document Tools — Generate PDF reports, PPTX slides, and CSV exports.

Uses fpdf2 for PDF, python-pptx for slides, csv module for CSV.
Generated files are saved to /workspace/output/ for Telegram delivery.
"""
from __future__ import annotations

import csv
import io
import json
import logging
import os
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional

logger = logging.getLogger(__name__)

OUTPUT_DIR = Path(os.getenv("TOKIO_OUTPUT_DIR", "/workspace/output"))


def _ensure_output_dir() -> Path:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    return OUTPUT_DIR


# ── PDF Generation ────────────────────────────────────────────────────────

def _sanitize_text(text: str) -> str:
    """Replace problematic Unicode chars with ASCII equivalents for PDF safety."""
    replacements = {
        "\u2014": "-",   # em dash
        "\u2013": "-",   # en dash
        "\u2018": "'",   # left single quote
        "\u2019": "'",   # right single quote
        "\u201c": '"',   # left double quote
        "\u201d": '"',   # right double quote
        "\u2026": "...", # ellipsis
        "\u2022": "*",   # bullet
        "\u00b7": "*",   # middle dot
        "\u2192": "->",  # right arrow
        "\u2190": "<-",  # left arrow
        "\u2264": "<=",  # less-equal
        "\u2265": ">=",  # greater-equal
        "\u2260": "!=",  # not-equal
    }
    for old, new in replacements.items():
        text = text.replace(old, new)
    return text


def _find_dejavu_font() -> Optional[str]:
    """Find DejaVuSans.ttf on the filesystem."""
    candidates = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/DejaVuSans.ttf",
        "/usr/share/fonts/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/DejaVuSans.ttf",
    ]
    for path in candidates:
        if os.path.isfile(path):
            return path
    # Search more broadly
    for root, _dirs, files in os.walk("/usr/share/fonts"):
        for f in files:
            if f == "DejaVuSans.ttf":
                return os.path.join(root, f)
    return None


def _find_dejavu_bold() -> Optional[str]:
    """Find DejaVuSans-Bold.ttf on the filesystem."""
    candidates = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
        "/usr/share/fonts/truetype/DejaVuSans-Bold.ttf",
        "/usr/share/fonts/dejavu/DejaVuSans-Bold.ttf",
    ]
    for path in candidates:
        if os.path.isfile(path):
            return path
    for root, _dirs, files in os.walk("/usr/share/fonts"):
        for f in files:
            if f == "DejaVuSans-Bold.ttf":
                return os.path.join(root, f)
    return None


def _generate_pdf(
    title: str,
    sections: List[Dict[str, str]],
    output_path: str = "",
    template: str = "default",
) -> str:
    """Generate a PDF report.

    Args:
        title: Report title.
        sections: List of {"heading": "...", "body": "..."} dicts.
        output_path: Where to save. Defaults to /workspace/output/<title>.pdf.
        template: Template style (default, security, infrastructure).
    """
    try:
        from fpdf import FPDF
    except ImportError:
        return json.dumps({"ok": False, "error": "fpdf2 no instalado. Ejecuta: pip install fpdf2"})

    if not output_path:
        safe_title = "".join(c if c.isalnum() or c in " _-" else "_" for c in title)[:60]
        output_path = str(_ensure_output_dir() / f"{safe_title}_{datetime.now().strftime('%Y%m%d_%H%M')}.pdf")

    # Color schemes by template
    colors = {
        "default": {"header": (41, 128, 185), "accent": (52, 152, 219)},
        "security": {"header": (192, 57, 43), "accent": (231, 76, 60)},
        "infrastructure": {"header": (39, 174, 96), "accent": (46, 204, 113)},
    }
    scheme = colors.get(template, colors["default"])

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)

    # Try to load DejaVu TTF for full Unicode support
    dejavu_regular = _find_dejavu_font()
    dejavu_bold = _find_dejavu_bold()
    use_unicode = False

    if dejavu_regular:
        try:
            pdf.add_font("DejaVu", "", dejavu_regular)
            if dejavu_bold:
                pdf.add_font("DejaVu", "B", dejavu_bold)
            else:
                pdf.add_font("DejaVu", "B", dejavu_regular)
            use_unicode = True
            font_family = "DejaVu"
            logger.info("Using DejaVu TTF font for Unicode PDF")
        except Exception as e:
            logger.warning(f"Failed to load DejaVu font: {e}, falling back to Helvetica")
            font_family = "Helvetica"
    else:
        font_family = "Helvetica"
        logger.info("DejaVu font not found, using Helvetica (ASCII only)")

    def safe(text: str) -> str:
        """Make text safe for the current font."""
        if use_unicode:
            return text
        return _sanitize_text(text)

    # Title page
    pdf.add_page()
    pdf.set_fill_color(*scheme["header"])
    pdf.rect(0, 0, 210, 50, "F")
    pdf.set_text_color(255, 255, 255)
    pdf.set_font(font_family, "B", 22)
    pdf.set_y(15)
    pdf.cell(0, 12, safe(title), align="C", new_x="LMARGIN", new_y="NEXT")
    pdf.set_font(font_family, "", 10)
    pdf.cell(0, 8, f"Generado: {datetime.now().strftime('%Y-%m-%d %H:%M')}", align="C",
             new_x="LMARGIN", new_y="NEXT")
    pdf.cell(0, 6, "TokioAI v2.1", align="C", new_x="LMARGIN", new_y="NEXT")

    pdf.set_text_color(0, 0, 0)
    pdf.ln(20)

    # Content sections
    for section in sections:
        heading = section.get("heading", "")
        body = section.get("body", "")

        if heading:
            pdf.set_fill_color(*scheme["accent"])
            pdf.set_text_color(255, 255, 255)
            pdf.set_font(font_family, "B", 13)
            pdf.cell(0, 9, safe(f"  {heading}"), fill=True, new_x="LMARGIN", new_y="NEXT")
            pdf.ln(3)

        if body:
            pdf.set_text_color(50, 50, 50)
            pdf.set_font(font_family, "", 10)
            for line in body.split("\n"):
                pdf.multi_cell(0, 5, safe(line))
            pdf.ln(5)

    # Footer
    pdf.set_y(-25)
    pdf.set_font(font_family, "", 8)
    pdf.set_text_color(150, 150, 150)
    pdf.cell(0, 5, "Generado por TokioAI - Agente Autonomo", align="C")

    # Save
    try:
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        pdf.output(output_path)
        size = os.path.getsize(output_path)
        return json.dumps({
            "ok": True,
            "file": output_path,
            "size_bytes": size,
            "pages": pdf.page,
        }, ensure_ascii=False)
    except Exception as e:
        return json.dumps({"ok": False, "error": f"Error guardando PDF: {e}"})


# ── Slides Generation ─────────────────────────────────────────────────────

def _generate_slides(
    title: str,
    slides: List[Dict[str, Any]],
    output_path: str = "",
    template: str = "default",
) -> str:
    """Generate PPTX slides.

    Args:
        title: Presentation title.
        slides: List of {"title": "...", "content": "...", "bullets": [...]} dicts.
        output_path: Where to save.
        template: Template style.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        return json.dumps({"ok": False, "error": "python-pptx no instalado. Ejecuta: pip install python-pptx"})

    if not output_path:
        safe_title = "".join(c if c.isalnum() or c in " _-" else "_" for c in title)[:60]
        output_path = str(_ensure_output_dir() / f"{safe_title}_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    subtitle = slide.placeholders[1]
    subtitle.text = f"Generado: {datetime.now().strftime('%Y-%m-%d %H:%M')}\nTokioAI v2.1"

    # Content slides
    for slide_data in slides:
        slide_title = slide_data.get("title", "")
        content = slide_data.get("content", "")
        bullets = slide_data.get("bullets", [])

        if bullets:
            layout = prs.slide_layouts[1]  # Title and Content
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = slide_title
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            for i, bullet in enumerate(bullets):
                if i == 0:
                    tf.text = str(bullet)
                else:
                    p = tf.add_paragraph()
                    p.text = str(bullet)
                    p.level = 0
        else:
            layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = slide_title
            body = slide.placeholders[1]
            body.text = content

    # Save
    try:
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_path)
        size = os.path.getsize(output_path)
        return json.dumps({
            "ok": True,
            "file": output_path,
            "size_bytes": size,
            "slides_count": len(prs.slides),
        }, ensure_ascii=False)
    except Exception as e:
        return json.dumps({"ok": False, "error": f"Error guardando PPTX: {e}"})


# ── CSV Generation ────────────────────────────────────────────────────────

def _generate_csv(
    data: List[List[Any]],
    output_path: str = "",
    headers: Optional[List[str]] = None,
) -> str:
    """Generate CSV file.

    Args:
        data: List of rows (each row is a list of values).
        output_path: Where to save.
        headers: Optional header row.
    """
    if not output_path:
        output_path = str(_ensure_output_dir() / f"export_{datetime.now().strftime('%Y%m%d_%H%M')}.csv")

    try:
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        with open(output_path, "w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            if headers:
                writer.writerow(headers)
            for row in data:
                writer.writerow(row)
        size = os.path.getsize(output_path)
        return json.dumps({
            "ok": True,
            "file": output_path,
            "size_bytes": size,
            "rows": len(data),
        }, ensure_ascii=False)
    except Exception as e:
        return json.dumps({"ok": False, "error": f"Error guardando CSV: {e}"})


# ── Unified Tool Entry Point ─────────────────────────────────────────────

async def document_tool(
    action: str,
    title: Optional[str] = None,
    content: Optional[str] = None,
    sections: Optional[List[Dict[str, str]]] = None,
    data: Optional[List[List[Any]]] = None,
    headers: Optional[List[str]] = None,
    output_path: Optional[str] = None,
    template: Optional[str] = None,
    params: Optional[Dict[str, Any]] = None,
) -> str:
    """
    Document generation tool.

    Actions:
      - generate_pdf: Create PDF report
      - generate_slides: Create PPTX
      - generate_csv: Export CSV
    """
    # Support legacy nested params format
    if params and isinstance(params, dict):
        title = title or params.get("title")
        content = content or params.get("content") or params.get("body") or params.get("text")
        sections = sections or params.get("sections")
        data = data or params.get("data")
        headers = headers or params.get("headers")
        output_path = output_path or params.get("output_path")
        template = template or params.get("template")

    action = (action or "").strip().lower()
    title = str(title or "Reporte TokioAI").strip()
    output_path = str(output_path or "").strip()
    template = str(template or "default").strip()

    try:
        if action == "generate_pdf":
            pdf_sections = []

            # Normalize sections
            if sections:
                if isinstance(sections, str):
                    pdf_sections = [{"heading": "", "body": sections}]
                elif isinstance(sections, list):
                    for s in sections:
                        if isinstance(s, str):
                            pdf_sections.append({"heading": "", "body": s})
                        elif isinstance(s, dict):
                            pdf_sections.append(s)

            # Fallback: use content as single body
            if not pdf_sections and content:
                if isinstance(content, list):
                    pdf_sections = [{"heading": "", "body": "\n".join(str(x) for x in content)}]
                else:
                    pdf_sections = [{"heading": "", "body": str(content)}]

            if not pdf_sections:
                return json.dumps({"ok": False, "error": (
                    "PDF vacio: necesito 'content' (texto) o 'sections' (lista de {heading, body}). "
                    "Ejemplo: document(action='generate_pdf', title='Mi Reporte', "
                    "content='Clima: 20C\\nSalud: OK')"
                )}, ensure_ascii=False)

            logger.info(f"PDF: title={title!r}, sections={len(pdf_sections)}")
            return _generate_pdf(title, pdf_sections, output_path, template)

        elif action == "generate_slides":
            slides = sections or []
            return _generate_slides(title, slides, output_path, template)

        elif action == "generate_csv":
            return _generate_csv(data or [], output_path, headers)

        return json.dumps({"ok": False, "error": f"Accion no soportada: {action}",
                          "supported": ["generate_pdf", "generate_slides", "generate_csv"]},
                         ensure_ascii=False)
    except Exception as e:
        return json.dumps({"ok": False, "action": action, "error": str(e)}, ensure_ascii=False)
